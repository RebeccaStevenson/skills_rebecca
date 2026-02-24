#!/usr/bin/env python3
"""Generate a PowerPoint presentation of figures, tables, and equations from a
Marker-converted paper folder.

Usage:
    python generate_paper_pptx.py <paper_folder> [--output <path.pptx>] [--skip-tables] [--skip-equations] [--experimental-figure-captions]

The paper folder should contain:
  - A primary .md file (the Marker conversion output)
  - A _citation.md file (enriched metadata)
  - Figure image files (_page_*_Figure_*.(jpeg|jpg|png|webp))
"""

import argparse
from io import BytesIO
import os
import re
import sys
from pathlib import Path

MIN_PYTHON = (3, 10)
PPTX_IMPORT_ERROR = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError as exc:
    PPTX_IMPORT_ERROR = exc
    Presentation = None

    # Fallback stubs keep parsing helpers importable for tests when python-pptx
    # is unavailable. Runtime generation still checks and exits with a clear error.
    def Inches(value):
        return value

    def Pt(value):
        return value

    def Emu(value):
        return value

    class RGBColor(tuple):
        def __new__(cls, r, g, b):
            return tuple.__new__(cls, (r, g, b))

    class _PPAlign:
        RIGHT = "RIGHT"
        CENTER = "CENTER"

    PP_ALIGN = _PPAlign()
    MSO_ANCHOR = None

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except ImportError:
    plt = None

# ---------------------------------------------------------------------------
# Slide dimensions (widescreen 16:9)
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color palette
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_COLOR = RGBColor(0x4A, 0x4A, 0x6A)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT_COLOR = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_ACCENT = RGBColor(0xE8, 0xF4, 0xF8)
SLIDE_NUM_COLOR = RGBColor(0x99, 0x99, 0x99)
CAPTION_LABEL_COLOR = RGBColor(0x2E, 0x86, 0xAB)

# Fonts
TITLE_FONT = "Helvetica"
BODY_FONT = "Helvetica"


# ---------------------------------------------------------------------------
# Parsing helpers
# ---------------------------------------------------------------------------

def find_primary_md(folder: Path) -> Path | None:
    """Find the primary markdown file (not _citation, not _no_refs)."""
    for f in sorted(folder.glob("*.md")):
        name = f.stem
        if name.endswith("_citation") or name.endswith("_no_refs"):
            continue
        return f
    return None


def find_citation_md(folder: Path) -> Path | None:
    """Find the _citation.md file."""
    for f in folder.glob("*_citation.md"):
        return f
    return None


def parse_citation(citation_path: Path) -> dict:
    """Extract metadata from _citation.md."""
    info = {"doi": "", "link": "", "citation": "", "abstract": ""}
    if not citation_path or not citation_path.exists():
        return info

    text = citation_path.read_text(encoding="utf-8")

    doi_m = re.search(r"\*\*DOI:\*\*\s*(.+)", text)
    if doi_m:
        info["doi"] = doi_m.group(1).strip()

    link_m = re.search(r"\*\*Link:\*\*\s*\[.*?\]\((.*?)\)", text)
    if link_m:
        info["link"] = link_m.group(1).strip()

    # APA citation: text between "## APA Citation" and the next "---"
    apa_m = re.search(r"## APA Citation\s*\n+(.+?)(?:\n---|\n## )", text, re.DOTALL)
    if apa_m:
        info["citation"] = apa_m.group(1).strip()

    abs_m = re.search(r"## Abstract\s*\n+(.+?)(?=\n##\s+|\n---\s*\n|\Z)", text, re.DOTALL)
    if abs_m:
        info["abstract"] = abs_m.group(1).strip()

    return info


def extract_authors_from_citation(citation: str) -> str:
    """Extract author names from an APA citation string.

    Example input: 'Anderson, K. L., Rajagovindan, R., ... (2009). Title. *Journal*...'
    Returns: 'Anderson, K. L., Rajagovindan, R., ...'
    """
    if not citation:
        return ""
    # APA format: authors precede the (Year) part
    m = re.match(r"(.+?)\s*\(\d{4}\)", citation)
    if m:
        return m.group(1).strip().rstrip(",").strip()
    return ""


def extract_title_from_citation(citation: str) -> str:
    """Extract the paper title from an APA-like citation string."""
    if not citation:
        return ""

    # Common APA form: "(Year). Title. *Journal*"
    m = re.search(r"\(\d{4}\)\.\s*(.+?)\.\s*\*", citation)
    if m:
        return m.group(1).strip()

    # Fallback: stop at the sentence boundary after year.
    m = re.search(r"\(\d{4}\)\.\s*(.+?)\.\s*[A-Z]", citation)
    if m:
        candidate = m.group(1).strip()
        if len(candidate.split()) >= 3:
            return candidate

    return ""


def is_likely_section_heading(text: str) -> bool:
    """Heuristic check to avoid using section headers as the paper title."""
    if not text:
        return False

    clean = re.sub(r"<[^>]+>", " ", text)
    clean = re.sub(r"[*_`~]+", "", clean)
    clean = re.sub(r"\s+", " ", clean).strip()
    if not clean:
        return False

    lowered = clean.lower()
    if re.match(r"^\d+(\.\d+)*\s+[a-z]", lowered):
        return True
    if lowered in {
        "article",
        "abstract",
        "introduction",
        "methods",
        "results",
        "discussion",
        "conclusion",
        "references",
        "contents",
    }:
        return True
    if lowered.startswith(("article ", "original article", "review article", "contents")):
        return True
    return False


def is_suspicious_title(title: str) -> bool:
    """Detect extracted titles that likely are not true paper titles."""
    if not title:
        return True

    clean = re.sub(r"<[^>]+>", " ", title)
    clean = re.sub(r"[*_`~]+", "", clean)
    clean = re.sub(r"\s+", " ", clean).strip()
    lowered = clean.lower()

    if not clean:
        return True
    if is_likely_section_heading(clean):
        return True
    if "open access article" in lowered:
        return True

    return False


def parse_title_authors(md_text: str, citation_info: dict | None = None) -> tuple[str, str]:
    """Extract title and authors from the markdown and/or citation metadata."""
    lines = md_text.split("\n")
    title = ""
    authors = ""

    # Find first H1 that looks like a paper title (skip journal name headers)
    h1_indices = [i for i, l in enumerate(lines) if l.startswith("# ") and not l.startswith("## ")]
    for idx in h1_indices:
        candidate = lines[idx][2:].strip()
        if is_likely_section_heading(candidate):
            continue
        # Skip very short headers that are likely journal names.
        if len(candidate) > 15 or idx > 0:
            title = candidate
            # Authors are typically in the lines following the title
            author_lines = []
            for j in range(idx + 1, min(idx + 10, len(lines))):
                line = lines[j].strip()
                if not line:
                    if author_lines:
                        break
                    continue
                if line.startswith("#"):
                    break
                # Stop at DOI, SUMMARY, or Keywords
                if "DOI" in line or line.startswith("# SUMMARY") or line.startswith("Keywords"):
                    break
                # Skip lines that look like abstracts (too long without superscripts)
                if len(line) > 200 and "<sup>" not in line:
                    break
                # Ignore image/link/boilerplate lines often emitted by PDF converters.
                if re.search(r"!\[.*?\]\([^)]+\)", line):
                    continue
                if line.startswith("<span") and line.endswith("</span>"):
                    continue
                if re.match(r"^(Received|Accepted|Published|Edited)\b", line, re.IGNORECASE):
                    continue
                if re.match(r"^(Figure|Table|Equation|Eq\.)\b", line, re.IGNORECASE):
                    continue
                author_lines.append(line)
            if author_lines:
                authors = " ".join(author_lines)
                # Clean up HTML tags
                authors = re.sub(r"<sup>.*?</sup>", "", authors)
                authors = re.sub(r"\s+", " ", authors).strip()
            break

    # Prefer citation-derived authors if markdown parsing produced
    # something suspicious (too long or clearly not author names).
    if citation_info:
        citation_authors = extract_authors_from_citation(citation_info.get("citation", ""))
        suspicious_authors = (
            not authors
            or len(authors) > 300
            or bool(re.search(r"!\[|_page_\d+_", authors))
            or bool(re.match(r"^(Received|Accepted|Published|Edited)\b", authors, re.IGNORECASE))
            or bool(re.match(r"^(Figure|Table|Equation|Eq\.)\b", authors, re.IGNORECASE))
        )
        if citation_authors and suspicious_authors:
            authors = citation_authors

    return title, authors


def normalize_md_line(line: str) -> str:
    """Normalize markdown-rich lines for pattern matching."""
    clean = line.strip()
    clean = re.sub(r"<[^>]+>", " ", clean)
    clean = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", clean)
    clean = re.sub(r"^\s*#{1,6}\s*", "", clean)
    clean = re.sub(r"[*_`~]+", "", clean)
    # Marker occasionally emits malformed bold around labels, e.g. "(A text".
    clean = re.sub(r"\(([A-Z])\s", r"(\1) ", clean)
    clean = re.sub(r"\s+", " ", clean).strip()
    return clean


def normalize_subcaption_line(line: str) -> str:
    """Normalize sub-caption labels like '(A) text' or 'A, text'."""
    clean = normalize_md_line(line)
    if not clean:
        return ""

    match = re.match(r"^\(([A-Za-z](?:-[A-Za-z])?)\)\s*(.+)$", clean)
    if not match:
        match = re.match(r"^([A-Za-z](?:-[A-Za-z])?)[)\].,:-]\s*(.+)$", clean)
    if not match:
        return ""

    label = match.group(1).upper()
    text = match.group(2).strip()
    if len(text) < 3:
        return ""
    if re.match(r"^(Figure|Table)\b", text, re.IGNORECASE):
        return ""

    return f"({label}) {text}"


def extract_inline_subcaptions(text: str) -> list[str]:
    """Extract inline sub-captions embedded in a long figure caption line."""
    if not text:
        return []

    label_pat = re.compile(r"\(([A-Za-z](?:-[A-Za-z])?)\)")
    all_matches = list(label_pat.finditer(text))
    if not all_matches:
        return []

    # Keep first occurrence of labels in forward alphabetical order to avoid
    # cross-reference fragments like "...same as (A)..." being treated as new items.
    kept = []
    seen_labels = set()
    last_ord = 0
    for m in all_matches:
        label = m.group(1).upper()
        parts = label.split("-", 1)
        start = parts[0][0]
        end = parts[-1][0]
        if start in seen_labels:
            continue
        cur_ord = ord(start)
        if cur_ord < last_ord:
            continue
        end_ord = ord(end)
        if end_ord < cur_ord:
            end_ord = cur_ord
        for c in range(cur_ord, end_ord + 1):
            seen_labels.add(chr(c))
        last_ord = end_ord
        kept.append(m)

    out = []
    for idx, m in enumerate(kept):
        end = kept[idx + 1].start() if idx + 1 < len(kept) else len(text)
        segment = text[m.start():end].strip(" .;")
        sub = normalize_subcaption_line(segment)
        if sub:
            out.append(sub)
    return out


def parse_figure_heading(line: str) -> tuple[str, list[str]] | None:
    """Parse a figure heading line and return title plus inline sub-captions."""
    clean = normalize_md_line(line)
    if not clean:
        return None

    match = re.match(
        r"^((?:Supplemental\s+)?(?:Figure|Fig\.?))\s+([A-Za-z0-9.\-]+)\s*[.:]?\s*(.*)$",
        clean,
        re.IGNORECASE,
    )
    if not match:
        return None

    fig_prefix = match.group(1).strip()
    fig_num = match.group(2).rstrip(".")
    remainder = match.group(3).strip()
    title = f"{fig_prefix} {fig_num}"
    if remainder:
        # Keep only the main title segment before (A)/(B)... blocks.
        lead = re.split(r"\([A-Za-z](?:-[A-Za-z])?\)", remainder, maxsplit=1)[0].strip(" .:-")
        if lead:
            title = f"{title}. {lead}"

    inline_subcaps = extract_inline_subcaptions(remainder)
    return title, inline_subcaps


def extract_figure_label(caption_title: str, fallback_index: int) -> str:
    """Extract a compact figure label like 'Figure 3' from caption heading text."""
    clean = normalize_md_line(caption_title)
    if clean:
        match = re.match(
            r"^(?:Supplemental\s+)?(?:Figure|Fig\.?)\s+([A-Za-z0-9.\-]+)",
            clean,
            re.IGNORECASE,
        )
        if match:
            figure_num = match.group(1).rstrip(".")
            return f"Figure {figure_num}"
    return f"Figure {fallback_index}"


def build_complete_figure_caption(caption_title: str, caption_body: list[str]) -> str:
    """Combine heading + body into one complete caption text block."""
    parts = []
    if caption_title:
        parts.append(caption_title.strip())
    for line in caption_body:
        line = line.strip()
        if line:
            parts.append(line)
    combined = " ".join(parts)
    return re.sub(r"\s+", " ", combined).strip()


def find_figure_caption(lines: list[str], image_idx: int) -> tuple[str, list[str]]:
    """Find the nearest caption heading and sub-captions for an image line."""
    heading_idx = None
    caption_title = ""
    caption_body = []

    # Preferred pattern: heading after image.
    for j in range(image_idx + 1, min(image_idx + 35, len(lines))):
        parsed = parse_figure_heading(lines[j])
        if parsed:
            heading_idx = j
            caption_title, caption_body = parsed
            break

    # Fallback: some files place heading before image.
    if heading_idx is None:
        for j in range(image_idx - 1, max(-1, image_idx - 15), -1):
            parsed = parse_figure_heading(lines[j])
            if parsed:
                heading_idx = j
                caption_title, caption_body = parsed
                break

    if heading_idx is None:
        return "", []

    # Gather nearby sub-captions after heading.
    blank_count = 0
    for k in range(heading_idx + 1, min(heading_idx + 28, len(lines))):
        if k == image_idx:
            continue

        raw = lines[k].strip()
        if not raw:
            blank_count += 1
            if blank_count >= 2 and caption_body:
                break
            continue
        blank_count = 0

        if raw.startswith("|"):
            break
        if re.search(r"!\[.*?\]\(", raw):
            break
        if parse_figure_heading(raw):
            break

        sub = normalize_subcaption_line(raw)
        if sub:
            caption_body.append(sub)
            continue

        clean = normalize_md_line(raw)
        if not clean:
            continue

        if caption_body and re.match(r"^[a-z]", clean):
            caption_body[-1] = f"{caption_body[-1]} {clean}".strip()
            continue
        if caption_body:
            break

    # De-duplicate while preserving order.
    deduped = []
    seen = set()
    for line in caption_body:
        if line in seen:
            continue
        seen.add(line)
        deduped.append(line)

    return caption_title, deduped


def extract_figures(md_text: str, folder: Path) -> list[dict]:
    """Extract figure slides from markdown, ignoring Marker Picture assets."""
    figures = []
    lines = md_text.split("\n")

    for i, line in enumerate(lines):
        # Match figure image references used by Marker:
        # ![](_page_X_Figure_Y.jpeg)
        # Use re.search because lines may start with <span> tags
        img_match = re.search(r"!\[.*?\]\(([^)]+)\)", line)
        if not img_match:
            continue

        img_file = img_match.group(1).strip()
        img_name = Path(img_file).name
        if not re.match(
            r"^_page_\d+_Figure_\d+\.(?:jpe?g|png|webp)$",
            img_name,
            re.IGNORECASE,
        ):
            continue

        img_path = folder / img_file
        if not img_path.exists():
            continue

        caption_title, caption_body = find_figure_caption(lines, i)
        figure_idx = len(figures) + 1
        figure_label = extract_figure_label(caption_title, figure_idx)
        caption_text = build_complete_figure_caption(caption_title, caption_body)
        if not caption_text:
            caption_text = figure_label

        figures.append({
            "type": "figure",
            "image_path": str(img_path),
            "label": figure_label,
            "caption_text": caption_text,
            "title": caption_title,
            "body": caption_body,
            "order": i,
        })

    return figures


def extract_tables(md_text: str) -> list[dict]:
    """Extract markdown tables with their titles."""
    tables = []
    lines = md_text.split("\n")
    i = 0

    while i < len(lines):
        line = lines[i].strip()

        # Look for table heading or Table N. text
        table_heading = re.match(
            r"(?:#{1,6}\s+)?(Table\s+\S+\.?\s+.+)", line
        )

        if table_heading:
            title = table_heading.group(1)
            title = re.sub(r"^#{1,6}\s+", "", title)

            # Collect the markdown table that follows
            table_lines = []
            for j in range(i + 1, len(lines)):
                tl = lines[j].strip()
                if not tl:
                    if table_lines:
                        break
                    continue
                if tl.startswith("|"):
                    table_lines.append(tl)
                elif table_lines:
                    break
            if table_lines:
                tables.append({
                    "type": "table",
                    "title": title,
                    "raw_lines": table_lines,
                    "order": i,
                })
                i = j
                continue

        # Also detect tables without explicit headings (standalone pipe tables)
        if line.startswith("|") and i > 0:
            # Check if there's a "Table" mention in the preceding few lines
            title = ""
            for k in range(max(0, i - 3), i):
                prev = lines[k].strip()
                tm = re.match(r"(?:#{1,6}\s+)?(Table\s+\S+\.?\s+.+)", prev)
                if tm:
                    title = tm.group(1)
                    title = re.sub(r"^#{1,6}\s+", "", title)
                    break

            if title:
                table_lines = []
                for j in range(i, len(lines)):
                    tl = lines[j].strip()
                    if tl.startswith("|"):
                        table_lines.append(tl)
                    elif table_lines and not tl:
                        break
                    elif table_lines:
                        break

                if table_lines and len(table_lines) >= 2:
                    tables.append({
                        "type": "table",
                        "title": title,
                        "raw_lines": table_lines,
                        "order": i,
                    })
                    i = j
                    continue

        i += 1

    return tables


def extract_equations(md_text: str) -> list[dict]:
    """Extract block equations ($$...$$) with any equation numbers."""
    equations = []
    lines = md_text.split("\n")
    i = 0

    def is_likely_equation_block(eq_text: str) -> bool:
        eq = eq_text.strip()
        if not eq:
            return False
        if len(eq) > 1600:
            return False

        eq_lines = [ln.strip() for ln in eq.splitlines() if ln.strip()]
        if not eq_lines or len(eq_lines) > 14:
            return False
        if any(ln.startswith("|") for ln in eq_lines):
            return False
        if re.search(r"https?://", eq):
            return False

        math_tokens = re.findall(r"(\\[A-Za-z]+|[=+\-*/^_{}]|≤|≥|∑|∫|≈|≠|±)", eq)
        if len(math_tokens) < 2:
            return False

        text_without_latex_cmds = re.sub(r"\\[A-Za-z]+", " ", eq)
        long_words = re.findall(r"[A-Za-z]{3,}", text_without_latex_cmds)
        sentence_marks = len(re.findall(r"[.!?]", eq))

        # Reject blocks that look mostly like prose rather than math.
        if len(long_words) > 45 and len(math_tokens) < 10:
            return False
        if sentence_marks >= 3 and len(math_tokens) < 12:
            return False

        return True

    while i < len(lines):
        if lines[i].strip() == "$$":
            eq_lines = []
            j = i + 1
            while j < len(lines) and lines[j].strip() != "$$":
                eq_lines.append(lines[j])
                j += 1

            if j < len(lines):
                eq_content = "\n".join(eq_lines).strip()
                if not is_likely_equation_block(eq_content):
                    i = j + 1
                    continue

                # Look for equation number on the next line
                eq_number = ""
                tag_match = re.search(r"\\tag\{([^}]+)\}", eq_content)
                if tag_match:
                    eq_number = tag_match.group(1).strip()

                if j + 1 < len(lines):
                    for off in (1, 2):
                        if j + off >= len(lines):
                            break
                        next_line = normalize_md_line(lines[j + off])
                        num_match = re.search(r"(?:\\n)?\(([A-Za-z]*\d+[A-Za-z.\-]*)\)", next_line)
                        if num_match:
                            eq_number = num_match.group(1)
                            break

                # Look for context: the paragraph before the equation
                context = ""
                for k in range(i - 1, max(0, i - 5) - 1, -1):
                    prev = normalize_md_line(lines[k])
                    if prev and not prev.startswith("#") and not prev.startswith("$$"):
                        # Truncate to last sentence
                        sentences = re.split(r"(?<=[.!?])\s+", prev)
                        if sentences:
                            context = sentences[-1][:200]
                        break

                label = f"Equation ({eq_number})" if eq_number else "Equation"
                equations.append({
                    "type": "equation",
                    "label": label,
                    "latex": eq_content,
                    "context": context,
                    "order": i,
                })
                i = j + 1
                continue
        i += 1

    return equations


def parse_table_to_rows(raw_lines: list[str]) -> list[list[str]]:
    """Parse markdown table lines into a list of rows (list of cells)."""
    rows = []
    for line in raw_lines:
        # Skip separator rows
        if re.match(r"^\|[\s\-:|]+\|$", line):
            continue
        cells = [c.strip() for c in line.split("|")]
        # Remove empty first/last from leading/trailing pipes
        if cells and cells[0] == "":
            cells = cells[1:]
        if cells and cells[-1] == "":
            cells = cells[:-1]
        if cells:
            rows.append(cells)
    return rows


def sanitize_latex_for_rendering(latex: str) -> list[str]:
    """Convert LaTeX block text into one or more renderable math lines."""
    if not latex:
        return []

    eq = latex.strip()
    eq = re.sub(r"\\tag\{[^}]*\}", "", eq)
    eq = re.sub(r"\\label\{[^}]*\}", "", eq)
    eq = re.sub(r"\\begin\{(?:aligned|align\*?|gather\*?|eqnarray\*?)\}", "", eq)
    eq = re.sub(r"\\end\{(?:aligned|align\*?|gather\*?|eqnarray\*?)\}", "", eq)
    eq = eq.replace("&", "")

    raw_lines = [ln.strip() for ln in re.split(r"\\\\|\n+", eq) if ln.strip()]
    lines = []
    for ln in raw_lines:
        ln = re.sub(r"^\$+|\$+$", "", ln).strip()
        if ln:
            lines.append(ln)
    return lines


def render_equation_png(latex: str) -> BytesIO | None:
    """Render LaTeX-ish math into a transparent PNG buffer using matplotlib."""
    if plt is None:
        return None

    lines = sanitize_latex_for_rendering(latex)
    if not lines:
        return None

    try:
        fig_h = max(1.2, min(3.2, 0.75 + 0.55 * len(lines)))
        fig = plt.figure(figsize=(10.0, fig_h), dpi=300)
        fig.patch.set_alpha(0.0)
        ax = fig.add_axes([0, 0, 1, 1])
        ax.axis("off")

        top_y = 0.5 + (len(lines) - 1) * 0.15
        for idx, line in enumerate(lines):
            y = top_y - idx * 0.30
            ax.text(
                0.5,
                y,
                rf"${line}$",
                ha="center",
                va="center",
                fontsize=30,
                color="#1A1A2E",
            )

        buf = BytesIO()
        fig.savefig(buf, format="png", transparent=True, bbox_inches="tight", pad_inches=0.08)
        plt.close(fig)
        buf.seek(0)
        return buf
    except Exception:
        try:
            plt.close(fig)
        except Exception:
            pass
        return None


# ---------------------------------------------------------------------------
# Slide creation helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color=BG_COLOR):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number(slide, num, total):
    """Add a small slide number in the bottom-right corner."""
    left = SLIDE_WIDTH - Inches(1.2)
    top = SLIDE_HEIGHT - Inches(0.45)
    txBox = slide.shapes.add_textbox(left, top, Inches(1.0), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = SLIDE_NUM_COLOR
    p.alignment = PP_ALIGN.RIGHT


def add_accent_bar(slide, left, top, width, height, color=ACCENT_COLOR):
    """Add a thin accent rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def create_title_slide(prs, title, authors, citation_info):
    """Create the title/metadata slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide)

    # Accent bar at top
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))

    # Title
    left, top = Inches(1.0), Inches(1.2)
    width, height = Inches(11.3), Inches(2.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.font.name = TITLE_FONT

    # Authors
    if authors:
        top = Inches(3.4)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = authors
        p.font.size = Pt(16)
        p.font.color.rgb = SUBTITLE_COLOR
        p.font.name = BODY_FONT

    # Citation
    citation = citation_info.get("citation", "")
    if citation:
        top = Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = citation
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.font.name = BODY_FONT

    # DOI
    doi = citation_info.get("doi", "")
    if doi:
        top = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"DOI: {doi}"
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_COLOR
        p.font.name = BODY_FONT

    # Abstract (truncated)
    abstract = citation_info.get("abstract", "")
    if abstract:
        top = Inches(6.0)
        txBox = slide.shapes.add_textbox(left, top, width, Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        truncated = abstract[:400] + ("..." if len(abstract) > 400 else "")
        p.text = truncated
        p.font.size = Pt(10)
        p.font.color.rgb = BODY_COLOR
        p.font.name = BODY_FONT

    # Bottom accent bar
    add_accent_bar(
        slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08)
    )


def get_image_dimensions(img_path: str, max_width: float, max_height: float):
    """Calculate image dimensions to fit within bounds while maintaining aspect ratio."""
    if PILImage:
        with PILImage.open(img_path) as img:
            img_w, img_h = img.size
    else:
        # Fallback: assume 4:3 aspect ratio
        img_w, img_h = 4, 3

    aspect = img_w / img_h

    # Fit within bounds
    if aspect > (max_width / max_height):
        # Width-constrained
        w = max_width
        h = max_width / aspect
    else:
        # Height-constrained
        h = max_height
        w = max_height * aspect

    return w, h


def create_figure_slide(
    prs,
    figure: dict,
    slide_num: int,
    total: int,
    include_caption_text: bool = False,
):
    """Create a slide for a figure.

    Default mode is figure-only. Caption text rendering is optional and
    considered experimental because OCR/markdown caption extraction can be noisy.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Top accent bar
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))

    # Figure title at top
    title_text = figure.get("label", "Figure")
    left, top = Inches(0.6), Inches(0.25)
    txBox = slide.shapes.add_textbox(left, top, Inches(12.0), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CAPTION_LABEL_COLOR
    p.font.name = TITLE_FONT

    # Layout defaults to figure-only. In experimental mode we reserve a
    # right caption column and render the extracted caption block.
    img_path = figure.get("image_path", "")
    caption_text = figure.get("caption_text", "").strip() if include_caption_text else ""

    img_left = Inches(0.5)
    img_top = Inches(1.1)
    content_height = Inches(5.8)

    if caption_text:
        cap_width = Inches(4.3)
        cap_left = SLIDE_WIDTH - cap_width - Inches(0.5)  # right-aligned with margin
        img_max_w = cap_left - img_left - Inches(0.3)  # 0.3" gap
    else:
        img_max_w = SLIDE_WIDTH - Inches(1.0)
    img_max_h = content_height

    if img_path and os.path.exists(img_path):
        w, h = get_image_dimensions(img_path, img_max_w, img_max_h)
        # Center image vertically in its area
        img_y_offset = (img_max_h - h) / 2
        slide.shapes.add_picture(
            img_path, img_left, img_top + img_y_offset, int(w), int(h)
        )

    # Full caption text on the right side (experimental mode only)
    if caption_text:
        cap_top = Inches(1.1)
        cap_height = content_height
        txBox = slide.shapes.add_textbox(cap_left, cap_top, cap_width, cap_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        clean = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", caption_text)
        clean = re.sub(r"<[^>]+>", "", clean)
        clean = re.sub(r"\s+", " ", clean).strip()
        p.text = clean
        p.font.size = Pt(11)
        p.font.color.rgb = BODY_COLOR
        p.font.name = BODY_FONT
        p.space_after = Pt(0)

    add_slide_number(slide, slide_num, total)


def create_table_slide(prs, table: dict, slide_num: int, total: int):
    """Create a slide for a markdown table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))

    # Title
    title_text = table.get("title", "Table")
    left, top = Inches(0.6), Inches(0.25)
    txBox = slide.shapes.add_textbox(left, top, Inches(12.0), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    tbl_match = re.match(r"(Table\s+\S+\.?)\s*(.*)", title_text)
    if tbl_match:
        run1 = p.add_run()
        run1.text = tbl_match.group(1) + "  "
        run1.font.size = Pt(18)
        run1.font.bold = True
        run1.font.color.rgb = CAPTION_LABEL_COLOR
        run1.font.name = TITLE_FONT

        run2 = p.add_run()
        run2.text = tbl_match.group(2)
        run2.font.size = Pt(18)
        run2.font.bold = True
        run2.font.color.rgb = TITLE_COLOR
        run2.font.name = TITLE_FONT
    else:
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.font.name = TITLE_FONT

    # Parse table
    rows = parse_table_to_rows(table.get("raw_lines", []))
    if not rows:
        return

    n_rows = min(len(rows), 20)  # Cap to avoid overflow
    n_cols = max(len(r) for r in rows)

    # Create PowerPoint table
    tbl_left = Inches(0.6)
    tbl_top = Inches(1.2)
    tbl_width = Inches(12.0)
    row_height = Inches(0.35)
    tbl_height = row_height * n_rows

    pptx_table = slide.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height
    ).table

    for r_idx in range(n_rows):
        row_data = rows[r_idx] if r_idx < len(rows) else []
        for c_idx in range(n_cols):
            cell = pptx_table.cell(r_idx, c_idx)
            cell_text = row_data[c_idx] if c_idx < len(row_data) else ""
            # Clean HTML
            cell_text = re.sub(r"<br\s*/?>", " ", cell_text)
            cell_text = re.sub(r"<[^>]+>", "", cell_text)
            cell.text = cell_text.strip()

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.name = BODY_FONT
                paragraph.font.color.rgb = BODY_COLOR

            # Header row styling
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_ACCENT
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(10)

    add_slide_number(slide, slide_num, total)


def create_equation_slide(prs, equation: dict, slide_num: int, total: int):
    """Create a slide for an equation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))

    # Label
    label = equation.get("label", "Equation")
    left, top = Inches(0.6), Inches(0.25)
    txBox = slide.shapes.add_textbox(left, top, Inches(12.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CAPTION_LABEL_COLOR
    p.font.name = TITLE_FONT

    # Equation in a centered box with light background
    latex = equation.get("latex", "")
    eq_left = Inches(1.5)
    eq_top = Inches(2.5)
    eq_width = Inches(10.3)
    eq_height = Inches(2.0)

    # Background box
    bg_shape = slide.shapes.add_shape(1, eq_left, eq_top, eq_width, eq_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = LIGHT_ACCENT
    bg_shape.line.fill.background()

    rendered = render_equation_png(latex)
    if rendered is not None:
        inner_pad = Inches(0.25)
        box_w = eq_width - inner_pad * 2
        box_h = eq_height - inner_pad * 2
        img_w, img_h = 4, 1
        if PILImage:
            with PILImage.open(rendered) as img:
                img_w, img_h = img.size
            rendered.seek(0)

        aspect = img_w / max(1, img_h)
        max_aspect = box_w / box_h
        if aspect >= max_aspect:
            draw_w = box_w
            draw_h = box_w / aspect
        else:
            draw_h = box_h
            draw_w = box_h * aspect

        draw_left = eq_left + (eq_width - draw_w) / 2
        draw_top = eq_top + (eq_height - draw_h) / 2
        slide.shapes.add_picture(rendered, int(draw_left), int(draw_top), int(draw_w), int(draw_h))
    else:
        # Fallback: plain LaTeX source when renderer/dependencies are unavailable.
        txBox = slide.shapes.add_textbox(
            eq_left + Inches(0.3), eq_top + Inches(0.3),
            eq_width - Inches(0.6), eq_height - Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = latex
        p.font.size = Pt(16)
        p.font.color.rgb = TITLE_COLOR
        p.font.name = "Courier New"
        p.alignment = PP_ALIGN.CENTER

    # Context paragraph
    context = equation.get("context", "")
    if context:
        ctx_top = Inches(5.0)
        txBox = slide.shapes.add_textbox(Inches(1.5), ctx_top, Inches(10.3), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        clean = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", context)
        clean = re.sub(r"<[^>]+>", "", clean)
        p.text = clean.strip()
        p.font.size = Pt(12)
        p.font.color.rgb = BODY_COLOR
        p.font.name = BODY_FONT

    add_slide_number(slide, slide_num, total)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def ensure_runtime_requirements():
    if sys.version_info < MIN_PYTHON:
        min_version = ".".join(str(v) for v in MIN_PYTHON)
        print(f"ERROR: Python {min_version}+ is required. Found {sys.version.split()[0]}.")
        sys.exit(1)

    if Presentation is None:
        print("ERROR: python-pptx is required. Install with: pip install python-pptx")
        if PPTX_IMPORT_ERROR:
            print(f"Import error: {PPTX_IMPORT_ERROR}")
        sys.exit(1)


def combine_items_in_document_order(*collections: list[dict]) -> list[dict]:
    """Sort heterogeneous extracted items by their markdown line index."""
    items = []
    for collection in collections:
        items.extend(collection)
    return sorted(items, key=lambda item: item.get("order", 10**9))


def generate_pptx(
    folder: str,
    output: str | None = None,
    skip_tables: bool = False,
    skip_equations: bool = False,
    experimental_figure_captions: bool = False,
):
    ensure_runtime_requirements()

    folder = Path(folder).resolve()
    if not folder.is_dir():
        print(f"ERROR: {folder} is not a directory")
        sys.exit(1)

    # Find files
    md_path = find_primary_md(folder)
    if not md_path:
        print(f"ERROR: No primary .md file found in {folder}")
        sys.exit(1)

    citation_path = find_citation_md(folder)

    print(f"Paper folder: {folder.name}")
    print(f"Markdown:     {md_path.name}")
    print(f"Citation:     {citation_path.name if citation_path else 'not found'}")

    # Parse
    md_text = md_path.read_text(encoding="utf-8")
    citation_info = parse_citation(citation_path)
    title, authors = parse_title_authors(md_text, citation_info)

    # If title parsing failed or looks suspicious, try from citation metadata.
    if is_suspicious_title(title):
        citation_title = extract_title_from_citation(citation_info.get("citation", ""))
        if citation_title:
            title = citation_title

    if not title:
        title = folder.name.replace("_", " ")

    print(f"Title:        {title[:80]}...")
    print(f"Authors:      {authors[:80]}..." if authors else "Authors:      (none found)")

    # Extract content
    figures = extract_figures(md_text, folder)
    tables = extract_tables(md_text) if not skip_tables else []
    equations = extract_equations(md_text) if not skip_equations else []

    print(f"Figures:      {len(figures)}")
    print(f"Tables:       {len(tables)}")
    print(f"Equations:    {len(equations)}")
    if experimental_figure_captions:
        print("Figure mode:  EXPERIMENTAL captions enabled")
    else:
        print("Figure mode:  Figure-only (default)")

    # Combine all items in markdown document order.
    all_items = combine_items_in_document_order(figures, tables, equations)

    if not all_items:
        print("WARNING: No figures, tables, or equations found. Creating title-only presentation.")

    total_slides = 1 + len(all_items)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    create_title_slide(prs, title, authors, citation_info)

    # Content slides
    for idx, item in enumerate(all_items):
        slide_num = idx + 2
        if item["type"] == "figure":
            create_figure_slide(
                prs,
                item,
                slide_num,
                total_slides,
                include_caption_text=experimental_figure_captions,
            )
        elif item["type"] == "table":
            create_table_slide(prs, item, slide_num, total_slides)
        elif item["type"] == "equation":
            create_equation_slide(prs, item, slide_num, total_slides)

    # Save
    if not output:
        output = str(folder / f"{folder.name}_figures.pptx")

    prs.save(output)
    print(f"\nSaved: {output}")
    print(f"Total slides: {total_slides}")
    return output


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Generate a PowerPoint of figures, tables, and equations from a Marker-converted paper."
    )
    parser.add_argument("paper_folder", help="Path to the paper folder in library_markdown_output/")
    parser.add_argument("--output", "-o", help="Output .pptx path (default: <folder>/<name>_figures.pptx)")
    parser.add_argument("--skip-tables", action="store_true", help="Skip table slides")
    parser.add_argument("--skip-equations", action="store_true", help="Skip equation slides")
    parser.add_argument(
        "--experimental-figure-captions",
        action="store_true",
        help="Render extracted figure captions on slides (experimental, prone to OCR/parse errors)",
    )
    parser.add_argument(
        "--with-captions",
        action="store_true",
        help="Alias for --experimental-figure-captions",
    )
    args = parser.parse_args()

    generate_pptx(
        args.paper_folder,
        output=args.output,
        skip_tables=args.skip_tables,
        skip_equations=args.skip_equations,
        experimental_figure_captions=(
            args.experimental_figure_captions or args.with_captions
        ),
    )
